import json
import uuid
import boto3
import constants
import utilities
from utilities import JsonHelper
from utilities import Helper
from utilities import Logger
from utilities import LogLevel
from botocore.exceptions import ClientError
from dynamoDBHelper import DynamoDBHelper
from appConfigHelper import AppConfigHelper
from dynamoDBHelper import DynamoDBHelper
from utilities import Helper
from datetime import datetime
from pptx import Presentation 
from pptx.util import Inches 


#Amazon Transcribe JSON schema constants
TRANSCRIBE_JSE__END_TIME = 'end_time'
TRANSCRIBE_JSE__ALTERNATIVES = 'alternatives'
TRANSCRIBE_JSE__PUNCTUATION = 'punctuation'
TRANSCRIBE_JSE__CONTENT = 'content'
TRANSCRIBE_JSE__TYPE = 'type'
TRANSCRIBE_JSE__RESULTS = 'results'
TRANSCRIBE_JSE__ITEMS = 'items'
    

class PresentationHelper:
   
    def __init__(self, slideTransistionFilename, transcriptionFilename):
        self.slideTransistionFilename = slideTransistionFilename
        self.transcriptionFilename = transcriptionFilename
    
    # Functions to merge slide transitions and transcriptions to create slide transciptions file
    def createNotes(slideNumber, slidePrefix, slideImageMimeType, slideTranscript):
        return json.dumps({
                          constants.UGH_PC_JSE__SLIDENUMBER : slideNumber + 1,
                          constants.UGH_PC_JSE__SLIDEIMAGE : slidePrefix + str(slideNumber + 1) + constants.CHAR__PERIOD + slideImageMimeType,
                          constants.UGH_PC_JSE__TRANSCRIPT : slideTranscript
                        })
    
    def getNumberOfSlides(presentationFilename: str) -> int:
        # Creating an Presentation object
        pptPresentation = Presentation(presentationFilename) 
        return totalPresentationSlides = len(pptPresentation.slides)
    
    
    def toJsonFile(notesAsJsonString, outputFilename = constants.CHAR__EMPTY):
        return utilities.JsonHelper.jsonStringToJsonFile(notesAsJsonString, outputFilename)
        
    
    def getNotes(self):
        slideTransistionFileData = utilities.JsonHelper.jsonFileToDictionary(self.slideTransistionFilename)
        transcriptionFileData = utilities.JsonHelper.jsonFileToDictionary(self.transcriptionFilename)
        Logger.log(f'GetNotes -> Successfully loaded data files 1. {constants.UGH_PC_DDB_FLD__AC_SLIDE_TRANSITIONS_FILENAME}={self.slideTransistionFilename}{constants.CHAR__NEWLINE} 2. {constants.UGH_PC_DDB_FLD__AC_TRANSCRIPTIONS_FILENAME}={self.transcriptionFilename}')

        if slideTransistionFileData is None:
            raise Exception(f"Error in reading the Slide-Transistions file : {slidetransistionfilename}")
        elif transcriptionFileData is None:
            raise Exception(f"Error in reading the Transcriptions file : {transcriptionfilename}")
        else:
            slideNumber = 0
            slideData = slideTransistionFileData[constants.UGH_PC_JSE__SLIDES][slideNumber]
            slideTranscript = constants.CHAR__EMPTY
            slideendtime = float(slideData[TRANSCRIBE_JSE__END_TIME])
            slideNotes = constants.CHAR__OPENCURLYBRACKET + constants.CHAR__SPACE + constants.CHAR__DOUBLE_QUOTE + constants.UGH_PC_JSE__SLIDES + constants.CHAR__DOUBLE_QUOTE + constants.CHAR__COLON + constants.CHAR__SPACE + constants.CHAR__OPENSQUAREBRACKET
            slideNotesDelimiter = constants.CHAR__EMPTY
            totalSlides = len(slideTransistionFileData[constants.UGH_PC_JSE__SLIDES])
            slidePrefix = slideTransistionFileData[constants.UGH_PC_JSE__SLIDEIMAGEPREFIX]
            slideImageMimeType = slideTransistionFileData[constants.UGH_PC_JSE__SLIDEIMAGEMIMETYPE]

            # Iterating through each word in the transcriptions data file
            for wordData in transcriptionFileData[TRANSCRIBE_JSE__RESULTS][TRANSCRIBE_JSE__ITEMS]:
                if wordData[TRANSCRIBE_JSE__TYPE] == TRANSCRIBE_JSE__PUNCTUATION:
                    slideTranscript = slideTranscript + wordData[TRANSCRIBE_JSE__ALTERNATIVES][0][TRANSCRIBE_JSE__CONTENT]            
                else:
                    wordendtime = float(wordData[TRANSCRIBE_JSE__END_TIME])
                    if wordendtime < slideendtime:
                        slideTranscript = slideTranscript + constants.CHAR__SPACE +  wordData[TRANSCRIBE_JSE__ALTERNATIVES][0][TRANSCRIBE_JSE__CONTENT]                
                    else:
                        slideNotes = slideNotes + slideNotesDelimiter + PresentationHelper.createNotes(slideNumber, slidePrefix, slideImageMimeType, slideTranscript)
                        slideNumber =  slideNumber + 1
                        slideData = slideTransistionFileData[constants.UGH_PC_JSE__SLIDES][slideNumber]
                        slideendtime = float(slideData[TRANSCRIBE_JSE__END_TIME])
                        slideTranscript = wordData[TRANSCRIBE_JSE__ALTERNATIVES][0][TRANSCRIBE_JSE__CONTENT]  
                        slideNotesDelimiter = constants.CHAR__COMMA                
            if slideNotes != constants.CHAR__EMPTY:
                if slideNumber + 1 == totalSlides:
                    slideNotes = slideNotes + slideNotesDelimiter + PresentationHelper.createNotes(slideNumber, slidePrefix, slideImageMimeType, slideTranscript)           
                slideNotes = slideNotes + constants.CHAR__CLOSESQUAREBRACKET + constants.CHAR__CLOSECURLYBRACKET
            Logger.log(f'GetNotes -> Successfully completely loading and merging slide transistions and transcriptions data. {constants.CHAR__NEWLINE} Result = {slideNotes}')
            return slideNotes

    # Functions to create/merge slide transciptions  as notes in to the powerpoint presentation file
    def getFormatedSlideNotes(overwriteNotes, slideTranscript, existingSlideTranscript):
        noteBeginLine = "\n\n#################################### BEGIN NOTES ######################################################\n"
        noteTimestampLine = f"-------------- NOTE ADDED ON: {Helper.getFormattedCurrentTime()} --------------------\n\n"
        noteEndLine =   "\n\n###################################### END NOTES ######################################################\n\n"
        return ('' if  overwriteNotes else existingSlideTranscript) + noteBeginLine + noteTimestampLine + slideTranscript + noteEndLine


    def addSlideImageToPresentation(presentationFilename, pptPresentation, imagesPath, slideNotes, overwriteNotes):
        
        Logger.log(f'AddSlideImageToPresentation -> presentationFilename={presentationFilename} {constants.CHAR__NEWLINE} imagesPath={imagesPath} {constants.CHAR__NEWLINE} slideNotes={slideNotes} {constants.CHAR__NEWLINE} overwriteNotes={overwriteNotes}')
        # Selecting blank slide
        blankSlideLayout = pptPresentation.slide_layouts[6] 

        # Attaching slide to ppt
        slide = pptPresentation.slides.add_slide(blankSlideLayout) 

        left = top = Inches(0) 
        picture = slide.shapes.add_picture(imagesPath, left, top, width=None, height=None )
        notesSlide = slide.notes_slide
        text_frame = notesSlide.notes_text_frame
        text_frame.text = PresentationHelper.getFormatedSlideNotes(overwriteNotes, slideNotes, '')
        # save file
        pptPresentation.save(presentationFilename)
        return


    def addSlideNotesToPresentation(presentationNotesParameters)-> bool:
        
        result = False
        Logger.log(f'AddSlideNotesToPresentation -> PresentationNotesParameters = {presentationNotesParameters}')
        
        if(presentationNotesParameters is None):
            raise Exception("[Presentation Notes Parameters] object cannot be None")

        slideTranscriptionsFilename = presentationNotesParameters[constants.UGH_PC_DDB_FLD__AC_SLIDE_TRANSCRIPTIONS_FILENAME] 
        presentationFilename = presentationNotesParameters[constants.UGH_PC_DDB_FLD__AC_PRESENTATION_FILENAME] 
        overwriteNotes = presentationNotesParameters[constants.UGH_PC_DDB_FLD__OVERWRITE_NOTES] 

        if presentationFilename is None or presentationFilename == constants.CHAR__EMPTY:
            if(presentationNotesParameters[constants.UGH_PC_DDB_FLD__AC_PRESENTATION_TEMPLATE_FILENAME] is None or presentationNotesParameters[constants.UGH_PC_DDB_FLD__AC_PRESENTATION_TEMPLATE_FILENAME] == ''):
                raise Exception("[presentation_template_filename] object cannot be None or Empty")    
            presentationFilename = presentationNotesParameters[constants.UGH_PC_DDB_FLD__AC_PRESENTATION_TEMPLATE_FILENAME]

        # Creating an Presentation object
        pptPresentation = Presentation(presentationFilename) 

        #Load slide transcriptions data from file
        slideTranscriptionFiledata = utilities.JsonHelper.jsonFileToDictionary(slideTranscriptionsFilename)
        totalSlideTranscriptionSlides = len(slideTranscriptionFiledata[constants.UGH_PC_JSE__SLIDES])
        totalPresentationSlides = len(pptPresentation.slides)

        if  totalSlideTranscriptionSlides <= totalPresentationSlides:
            Logger.log(f'AddSlideNotesToPresentation -> Slide count for [{slideTranscriptionsFilename}({totalSlideTranscriptionSlides})] is lesser than or equal to the source [{presentationFilename}({totalPresentationSlides})]. Using existing presentation.')
            slideIndex = 1
            for slide in pptPresentation.slides:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = PresentationHelper.getFormatedSlideNotes(overwriteNotes, slideTranscriptionFiledata[constants.UGH_PC_JSE__SLIDES][slideIndex-1][constants.UGH_PC_JSE__TRANSCRIPT], text_frame.text)
                Logger.log(f'AddSlideNotesToPresentation -> Adding notes to Slide # [{slideIndex} of {totalSlideTranscriptionSlides}]{constants.CHAR__NEWLINE} NOTE = [{text_frame.text}].')
                pptPresentation.save(presentationFilename)
                if slideIndex >= totalSlideTranscriptionSlides:
                    break
                slideIndex = slideIndex + 1            
        elif totalSlideTranscriptionSlides > totalPresentationSlides:
            Logger.log(f'AddSlideNotesToPresentation -> Slide count for [{slideTranscriptionsFilename}({totalSlideTranscriptionSlides})] is greater than the source [{presentationFilename}({totalPresentationSlides})]. Using images to recreate the presentation.')
            # Iterating through slide in the slide transcriptions data file
            for slide in slideTranscriptionFiledata[constants.UGH_PC_JSE__SLIDES]:
                # Giving Image path 
                imagesPath = presentationNotesParameters[constants.UGH_PC_DDB_FLD__AC_SLIDES_PREFIX] + slide[constants.UGH_PC_JSE__SLIDEIMAGE]
                PresentationHelper.addSlideImageToPresentation(presentationFilename, pptPresentation, imagesPath, slide[constants.UGH_PC_JSE__TRANSCRIPT], overwriteNotes)
        result = True
        return result